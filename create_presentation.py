"""
Generate ARIA Life feature deck → public/ARIA_Life_Features.pptx
Dark theme · #060608 background · electric blue + purple accents · clean minimal
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── Brand palette ────────────────────────────────────────────
BG          = RGBColor(6, 6, 8)          # #060608
BLUE        = RGBColor(59, 130, 246)     # #3b82f6
PURPLE      = RGBColor(124, 58, 237)     # #7c3aed
WHITE       = RGBColor(255, 255, 255)
TEXT_MUTED  = RGBColor(160, 160, 168)    # #a0a0a8
DIM         = RGBColor(74, 74, 85)       # #4a4a55
CARD        = RGBColor(15, 15, 18)       # #0f0f12
BORDER      = RGBColor(26, 26, 31)       # #1a1a1f

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)


def add_solid_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_accent_dot(slide, x, y, color=BLUE, size=Inches(0.18)):
    """Small colored circle used as a section marker."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=20, color=WHITE,
                bullet_color=BLUE, font='Calibri', spacing=14):
    """Custom bullet list — colored dot + text, max 5 items."""
    line_h = Inches(0.55)
    for i, item in enumerate(bullets[:5]):
        cy = y + line_h * i
        # Dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, cy + Inches(0.18),
            Inches(0.13), Inches(0.13),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.fill.background()
        # Text
        tb = slide.shapes.add_textbox(
            x + Inches(0.32), cy,
            w - Inches(0.32), line_h,
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_footer(slide, num, total):
    add_textbox(
        slide, Inches(0.6), Inches(7.05), Inches(4), Inches(0.3),
        'ARIA Life · arialife.app', size=10, color=DIM, font='Calibri',
    )
    add_textbox(
        slide, Inches(11.7), Inches(7.05), Inches(1.2), Inches(0.3),
        f'{num} / {total}', size=10, color=DIM, align=PP_ALIGN.RIGHT, font='Calibri',
    )


def title_slide(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_solid_bg(slide)

    # Subtle gradient accent strip on the left edge — vertical blue bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    # Big eyebrow
    add_textbox(
        slide, Inches(1.2), Inches(2.5), Inches(11), Inches(0.4),
        'ARIA LIFE', size=14, bold=True, color=BLUE, font='Calibri',
    )

    # Title
    add_textbox(
        slide, Inches(1.2), Inches(3.0), Inches(11), Inches(2.0),
        'AI Assistant for Serious Professionals',
        size=54, bold=True, color=WHITE, font='Calibri',
    )

    # Subtitle / URL
    add_textbox(
        slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5),
        'arialife.app',
        size=20, color=TEXT_MUTED, font='Calibri',
    )

    # Glow circle accents (right side)
    glow1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10.2), Inches(1.4), Inches(2.4), Inches(2.4),
    )
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = BLUE
    glow1.fill.transparency = 0.85
    glow1.line.fill.background()

    glow2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.0), Inches(2.5), Inches(1.6), Inches(1.6),
    )
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = PURPLE
    glow2.fill.transparency = 0.8
    glow2.line.fill.background()

    add_footer(slide, num, total)


def content_slide(prs, num, total, eyebrow, title, bullets, accent=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Eyebrow (slide number / category)
    add_textbox(
        slide, Inches(0.9), Inches(0.7), Inches(11), Inches(0.4),
        eyebrow, size=12, bold=True, color=accent, font='Calibri',
    )

    # Title
    add_textbox(
        slide, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.4),
        title, size=40, bold=True, color=WHITE, font='Calibri',
    )

    # Thin divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9), Inches(2.65),
        Inches(0.8), Emu(12700),  # 1pt height
    )
    div.fill.solid()
    div.fill.fore_color.rgb = accent
    div.line.fill.background()

    # Bullets
    add_bullets(
        slide, Inches(0.9), Inches(3.05),
        Inches(11.5), Inches(3.5),
        bullets, size=22, color=WHITE, bullet_color=accent,
    )

    add_footer(slide, num, total)


# ── Build deck ────────────────────────────────────────────────
SLIDES = [
    {
        'eyebrow': 'OVERVIEW',
        'title': 'What is ARIA Life',
        'bullets': [
            'AI business assistant built for serious professionals',
            'Handles emails, meetings, scheduling, and strategy',
            'Voice-powered — speak naturally, no typing required',
            'Available in English, Somali, Arabic, and Swahili',
            'Web app at arialife.app · iOS and Android coming',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'FEATURE 01',
        'title': 'Voice Planning',
        'bullets': [
            'Speak naturally to schedule meetings and tasks',
            'No typing, no clicking — just talk',
            'ARIA understands context and adds events instantly',
            'WhatsApp reminder sent automatically',
            'Works in all four supported languages',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'FEATURE 02',
        'title': 'Email AI Summaries',
        'bullets': [
            'Connects securely to your Gmail inbox',
            'Summarizes every email in one clear sentence',
            '50 emails processed in seconds',
            'Surfaces who sent it, what they need, and the action required',
            'Read the highlights, skip the noise',
        ],
        'accent': PURPLE,
    },
    {
        'eyebrow': 'FEATURE 03',
        'title': 'Meeting Recorder',
        'bullets': [
            'One-tap recording for any meeting',
            'Automatic transcription via AssemblyAI',
            'AI-generated summary with key decisions',
            'Action items extracted and assignable',
            'Searchable history of every meeting',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'FEATURE 04',
        'title': 'AI Strategy Chat',
        'bullets': [
            'Type or speak any business question',
            'Powered by Claude AI from Anthropic',
            'Expert answers in seconds, like a 24/7 consultant',
            'Save and export every conversation',
            'Follow-up questions remember full context',
        ],
        'accent': PURPLE,
    },
    {
        'eyebrow': 'STRATEGY EXAMPLES',
        'title': 'What You Can Ask',
        'bullets': [
            'Write a 12-month business plan for my startup',
            'Build a marketing strategy targeting SMEs',
            'Run a competitor analysis on three rivals',
            'Project monthly revenue and operating costs',
            'Draft a professional email for a difficult client',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'HOW IT WORKS',
        'title': 'AI Strategy in Four Steps',
        'bullets': [
            'You ask — type or speak a business question',
            'ARIA forwards your question to Claude AI',
            'Claude analyzes context and crafts a tailored answer',
            'Response is displayed instantly inside the app',
            'Follow-up questions supported in the same thread',
        ],
        'accent': PURPLE,
    },
    {
        'eyebrow': 'FEATURE 05',
        'title': 'WhatsApp Reminders',
        'bullets': [
            'Reminders sent directly to your WhatsApp number',
            'No new app to download — uses what you already have',
            'Powered by WasenderAPI for reliable delivery',
            'Works with regular WhatsApp, no Business account needed',
            'Available on Corporate Mini and above',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'FEATURE 06',
        'title': 'Smart Calendar',
        'bullets': [
            'Built-in calendar with manual scheduling',
            'Two-way Google Calendar sync',
            'Add events by voice or by tap',
            'Unified view of all events from every source',
            'Color-coded by type and source',
        ],
        'accent': PURPLE,
    },
    {
        'eyebrow': 'FEATURE 07',
        'title': 'Multi-Language Support',
        'bullets': [
            'English — primary interface',
            'Somali — full UI and voice translation',
            'Arabic — right-to-left layout supported',
            'Swahili — full UI and voice translation',
            'Switch instantly in Settings',
        ],
        'accent': BLUE,
    },
    {
        'eyebrow': 'GET STARTED',
        'title': 'Try ARIA Life Today',
        'bullets': [
            'Visit arialife.app',
            'Sign in with Google in seconds',
            '7-day free trial on every plan',
            'No credit card required',
            'Cancel anytime',
        ],
        'accent': PURPLE,
    },
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDES) + 1  # +1 for title slide
    title_slide(prs, 1, total)

    for i, s in enumerate(SLIDES, start=2):
        content_slide(
            prs, i, total,
            eyebrow=s['eyebrow'], title=s['title'],
            bullets=s['bullets'], accent=s['accent'],
        )

    out_dir = Path(__file__).parent / 'public'
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / 'ARIA_Life_Features.pptx'
    prs.save(out_path)
    print(f'✔ Generated {out_path}')


if __name__ == '__main__':
    main()
